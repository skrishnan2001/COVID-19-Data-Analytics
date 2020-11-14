from flask import Flask,redirect,url_for,render_template,request,Response
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib import style
import matplotlib.ticker as mtick
import numpy as np
from matplotlib.ticker import PercentFormatter
from matplotlib.figure import Figure
from matplotlib.backends.backend_agg import FigureCanvasAgg as FigureCanvas
import io
from pptx import Presentation
from pptx import util
import os
from District_Wise import main_func
from StatewiseGraphs import super_function
app = Flask(__name__)

@app.route('/')
def test():
    super_function()
    ppt()
    return "Successfully generated India-states graphs and its ppt."


def ppt():
    prs = Presentation()
    file_arr = []
    for file in os.listdir(os.path.join(os.getcwd(), 'images')):
        if file.endswith('.png'):
            file_arr.append(os.path.join(os.getcwd(),'images',file))
  
    for i in range(len(file_arr)):
         create_slide(file_arr[i],prs)
		
    prs.save('All_Graphs.pptx')

def create_slide(graph_name,prs):
   slide_layout = prs.slide_layouts[6]
   slide = prs.slides.add_slide(slide_layout)
   slide.shapes.add_picture(graph_name, 
                   left=util.Inches(0.1),
                   top=util.Inches(1.0),
                   width=util.Inches(9),
                   height=util.Inches(6))



@app.route('/tn')
def testTN():
    main_func()
    ppt1()
    return "Successfully generated TN-districts graphs and ppt"


def ppt1():
    prs = Presentation()
    file_arr = []
    for file in os.listdir(os.path.join(os.getcwd(), 'imagesTN')):
        if file.endswith('.png'):
            file_arr.append(os.path.join(os.getcwd(),'imagesTN',file))
  
    for i in range(len(file_arr)):
         create_slide1(file_arr[i],prs)
		
    prs.save('All_Graphs_TN.pptx')

def create_slide1(graph_name,prs):
   slide_layout = prs.slide_layouts[6]
   slide = prs.slides.add_slide(slide_layout)
   slide.shapes.add_picture(graph_name, 
                   left=util.Inches(0.1),
                   top=util.Inches(1.0),
                   width=util.Inches(9),
                   height=util.Inches(6))


if __name__ == "__main__":
    app.run(debug=True)
